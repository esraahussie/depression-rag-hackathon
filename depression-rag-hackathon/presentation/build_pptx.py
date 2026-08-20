"""Generate the MindCare 5-minute hackathon pitch deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap, qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

# ── Brand ──────────────────────────────────────────────────────────────────
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
TEAL_DEEP = RGBColor(0x11, 0x4B, 0x4A)
NAVY = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
SOFT = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF4, 0xF8, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x05, 0x96, 0x69)
CYAN = RGBColor(0x08, 0x91, 0xB2)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Calibri"
FONT_AR = "Segoe UI"

OUT = Path(__file__).resolve().parent / "MindCare_Hackathon_Pitch.pptx"


def set_run(run, text, size, color, bold=False, font=FONT, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.append(ea)
    ea.set("typeface", FONT_AR)


def add_textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
                font=FONT, italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, color, bold, font, italic)
    return box


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    set_run(p.add_run(), text, size, color, bold, italic=italic)
    return p


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def round_rect(slide, l, t, w, h, fill, line=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.25)
    s.shadow.inherit = False
    return s


def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def chevron(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, n, total=9, dark=False):
    color = RGBColor(0x9C, 0xA3, 0xAF) if dark else SOFT
    rect(slide, Inches(0), Inches(7.28), W, Inches(0.22), TEAL_DARK if dark else OFF_WHITE)
    add_textbox(slide, Inches(0.45), Inches(7.28), Inches(8), Inches(0.22),
                "MindCare  ·  Evidence-grounded clinical RAG  ·  ODC Hackathon",
                9, WHITE if dark else MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(11.4), Inches(7.28), Inches(1.5), Inches(0.22),
                f"{n}  /  {total}", 9, WHITE if dark else MUTED,
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def kicker(slide, l, t, text, color=TEAL):
    oval(slide, l, t + Inches(0.07), Inches(0.14), Inches(0.14), color)
    add_textbox(slide, l + Inches(0.22), t, Inches(8), Inches(0.28),
                text.upper(), 11, color, bold=True)


def card_title_body(slide, l, t, w, h, title, body, accent=TEAL, title_size=15, body_size=13):
    round_rect(slide, l, t, w, h, WHITE, LINE)
    rect(slide, l, t, Inches(0.08), h, accent)
    add_textbox(slide, l + Inches(0.22), t + Inches(0.12), w - Inches(0.34), Inches(0.36),
                title, title_size, SLATE, bold=True)
    add_textbox(slide, l + Inches(0.22), t + Inches(0.46), w - Inches(0.34), h - Inches(0.56),
                body, body_size, MUTED)


def numbered_badge(slide, l, t, n, fill=TEAL):
    oval(slide, l, t, Inches(0.36), Inches(0.36), fill)
    add_textbox(slide, l, t, Inches(0.36), Inches(0.36), str(n), 13, WHITE,
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── Slides ─────────────────────────────────────────────────────────────────

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.18), H, TEAL)
    rect(s, Inches(0), Inches(6.55), W, Inches(0.95), TEAL_DEEP)

    add_textbox(s, Inches(0.7), Inches(0.45), Inches(10), Inches(0.35),
                "ODC  ·  MENTAL HEALTH  ·  CLINICAL RAG HACKATHON", 13, TEAL, bold=True)
    add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(1.2),
                "MindCare", 54, WHITE, bold=True)
    add_textbox(s, Inches(0.7), Inches(2.5), Inches(11.5), Inches(0.9),
                "An evidence-grounded clinical assistant for depression.\nAnswers only from official guidelines — never from guesswork.",
                22, RGBColor(0xCB, 0xD5, 0xE1))
    add_textbox(s, Inches(0.7), Inches(3.7), Inches(11), Inches(0.4),
                "مساعد طبي مبني على الأدلة — بيرد من إرشادات NICE و WHO و USPSTF فقط",
                16, LIGHT_TEAL, font=FONT_AR)

    # three pills
    labels = [
        (Inches(0.7), "12 official PDFs"),
        (Inches(3.55), "848 cited chunks"),
        (Inches(6.55), "Zero-hallucination design"),
    ]
    for x, label in labels:
        round_rect(s, x, Inches(4.45), Inches(2.65), Inches(0.48), TEAL_DEEP, TEAL)
        add_textbox(s, x, Inches(4.45), Inches(2.65), Inches(0.48),
                    label, 14, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(s, Inches(0.7), Inches(6.68), Inches(8), Inches(0.55),
                "5-minute pitch  →  live demo  →  discussion     ·     We speak as a team",
                14, WHITE)
    add_textbox(s, Inches(10.2), Inches(6.68), Inches(2.7), Inches(0.55),
                "Slide 1 / 9", 13, LIGHT_TEAL, align=PP_ALIGN.RIGHT)

    notes(s, (
        "SPEAKER A — 10 seconds.\n"
        "We built MindCare: a clinical RAG assistant for depression. "
        "It answers only from official NICE, WHO, and USPSTF guidelines — "
        "and it is designed so the LLM cannot invent medical facts.\n"
        "احنا عملنا MindCare: مساعد طبي بيرد من الإرشادات الرسمية بس، ومن غير تخمين."
    ))


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "01  ·  Problem definition  ·  ~1 minute")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12), Inches(0.55),
                "Clinicians cannot safely ask a general LLM about depression.",
                26, NAVY, bold=True)

    # left problem cards
    card_title_body(
        s, Inches(0.5), Inches(1.35), Inches(6.0), Inches(1.55),
        "Guidelines exist — they are not usable at the bedside",
        "NICE, WHO and USPSTF publish hundreds of pages. A doctor, nurse, or caregiver cannot search them in seconds while a patient is waiting.",
        AMBER, 15, 14,
    )
    card_title_body(
        s, Inches(0.5), Inches(3.05), Inches(6.0), Inches(1.55),
        "Generic AI invents medical facts",
        "ChatGPT-style models guess dosages, mix adult and child pathways, and cite pages that do not exist. In mental health, hallucination is a safety failure — not a UX bug.",
        RED, 15, 14,
    )
    card_title_body(
        s, Inches(0.5), Inches(4.75), Inches(6.0), Inches(1.7),
        "Language and access gap in Egypt",
        "Most official evidence is English. Patients and junior staff need Egyptian Arabic, screening tools, and a refusal when the question is outside clinical scope.",
        CYAN, 15, 14,
    )

    # right: value
    round_rect(s, Inches(6.75), Inches(1.35), Inches(6.1), Inches(5.1), WHITE, LINE)
    rect(s, Inches(6.75), Inches(1.35), Inches(6.1), Inches(0.62), TEAL)
    add_textbox(s, Inches(6.9), Inches(1.42), Inches(5.8), Inches(0.5),
                "What we built — medical + operational value", 16, WHITE, bold=True)

    bullets = [
        ("Medical", "Every answer is grounded in retrieved guideline text with page-level citations."),
        ("Safety", "If evidence is weak or the question is off-topic, MindCare refuses to answer."),
        ("Clinical workflow", "PHQ-9, GAD-7 and EPDS screening → then ask the guideline about that score."),
        ("Operations", "Bilingual UI (English + المصري), voice in/out, FastAPI + React, ready to demo."),
        ("Market", "A decision-support layer for clinics, hotlines, and digital mental-health programs — not a chatbot that plays doctor."),
    ]
    y = 2.15
    for title, body in bullets:
        numbered_style = TEAL
        oval(s, Inches(7.0), Inches(y + 0.08), Inches(0.16), Inches(0.16), numbered_style)
        add_textbox(s, Inches(7.3), Inches(y), Inches(5.3), Inches(0.28), title, 14, TEAL_DARK, bold=True)
        add_textbox(s, Inches(7.3), Inches(y + 0.26), Inches(5.3), Inches(0.55), body, 13, MUTED)
        y += 0.82

    footer(s, 2)
    notes(s, (
        "SPEAKER A — 50–60 seconds. Speak as 'we'.\n\n"
        "The problem in one sentence: depression guidelines are official, long, and English — "
        "and a general LLM will invent an answer if you ask it anyway.\n\n"
        "Medical value: safer, cited answers from NICE / WHO / USPSTF. "
        "Operational value: Arabic, screening tools, and a refuse-to-answer path so the system never bluffs.\n\n"
        "This is a clinical decision-support product, not a therapy chatbot.\n\n"
        "بالعربي: المشكلة إن الإرشادات الرسمية طويلة، والـ LLM العادي بيختلق معلومات طبية. "
        "احنا بنحل ده بمساعد بيرد من الدليل بس، ولو مش متأكد بيرفض الإجابة."
    ))


def slide_ingestion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "02  ·  Data ingestion & processing")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12), Inches(0.5),
                "We start from official PDFs — not from the open web.",
                26, NAVY, bold=True)

    # source row
    sources = [
        ("NICE", "NG222 · NG134\nQS8 · QS48 · CG91\nAdults, children, chronic illness"),
        ("WHO", "EMRO clinical guides\nDepression factsheets\nRegional mental-health evidence"),
        ("USPSTF", "Depression & suicide screening\nAnxiety · perinatal\nAdults and young people"),
    ]
    x = Inches(0.5)
    for title, body in sources:
        round_rect(s, x, Inches(1.22), Inches(4.0), Inches(1.55), WHITE, LINE)
        rect(s, x, Inches(1.22), Inches(4.0), Inches(0.08), TEAL)
        add_textbox(s, x + Inches(0.2), Inches(1.38), Inches(3.6), Inches(0.32), title, 16, TEAL_DARK, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(1.72), Inches(3.6), Inches(0.9), body, 13, MUTED)
        x += Inches(4.15)

    # pipeline chevrons
    steps = [
        ("1. Extract", "pdfplumber\npage by page"),
        ("2. Clean", "headers, footers\nURLs, boilerplate"),
        ("3. Chunk", "1000 chars\n150 overlap"),
        ("4. Metadata", "PDF · page · chunk\npublic source URL"),
        ("5. Index", "MiniLM vectors\n+ BM25 keywords"),
    ]
    x = Inches(0.4)
    for i, (title, body) in enumerate(steps):
        if i < len(steps) - 1:
            chevron(s, x, Inches(3.05), Inches(2.55), Inches(1.15), TEAL if i % 2 == 0 else TEAL_DARK)
        else:
            round_rect(s, x, Inches(3.05), Inches(2.35), Inches(1.15), TEAL_DEEP)
        add_textbox(s, x + Inches(0.12), Inches(3.12), Inches(2.15), Inches(0.35), title, 14, WHITE, bold=True)
        add_textbox(s, x + Inches(0.12), Inches(3.46), Inches(2.15), Inches(0.65), body, 12, LIGHT_TEAL)
        x += Inches(2.55)

    # stats
    stats = [
        ("12", "official PDFs"),
        ("848", "evidence chunks"),
        ("page + URL", "on every chunk"),
        ("cached", "index rebuilds once"),
    ]
    x = Inches(0.5)
    for num, label in stats:
        round_rect(s, x, Inches(4.5), Inches(3.0), Inches(1.95), WHITE, LINE)
        add_textbox(s, x + Inches(0.15), Inches(4.7), Inches(2.7), Inches(0.7), num, 26, TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(5.45), Inches(2.7), Inches(0.7), label, 14, MUTED, align=PP_ALIGN.CENTER)
        x += Inches(3.2)

    footer(s, 3)
    notes(s, (
        "SPEAKER A — 45 seconds.\n\n"
        "We did not scrape the internet. We ingested 12 official PDFs: NICE (including adult treatment NG222 and children NG134), "
        "WHO EMRO, and USPSTF screening recommendations.\n\n"
        "Pipeline: extract page text → clean NICE footers, URLs, citation noise → recursive chunking "
        "(1000 characters, 150 overlap, drop tiny headers) → attach metadata: source file, page number, chunk index, public URL.\n\n"
        "Result: 848 chunks in Chroma + BM25, cached so the demo starts in seconds.\n\n"
        "بالعربي: بدأنا من الـ PDF الرسمي، نضفناه، قطعناه chunks، وكل chunk ماشي معاه اسم الملف ورقم الصفحة واللينك."
    ))


def slide_retrieval(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "03  ·  Retrieval & evidence layer")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12.3), Inches(0.5),
                "The LLM never searches. Retrieval finds the evidence first.",
                24, NAVY, bold=True)

    # user query bar
    round_rect(s, Inches(0.5), Inches(1.2), Inches(12.35), Inches(0.55), WHITE, TEAL)
    add_textbox(s, Inches(0.7), Inches(1.26), Inches(12), Inches(0.42),
                "User question  →  (Arabic is translated for search)  →  hybrid retriever  →  ranked evidence chunks",
                15, SLATE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    # two branches
    round_rect(s, Inches(0.5), Inches(2.0), Inches(5.9), Inches(2.15), WHITE, LINE)
    rect(s, Inches(0.5), Inches(2.0), Inches(5.9), Inches(0.48), TEAL)
    add_textbox(s, Inches(0.7), Inches(2.05), Inches(5.5), Inches(0.4), "Keyword path  ·  BM25", 16, WHITE, bold=True)
    add_textbox(s, Inches(0.7), Inches(2.6), Inches(5.5), Inches(1.4),
                "Catches exact clinical terms: PHQ-9, fluoxetine, NG222, lithium, screening.\nTop 50 candidates from the lexical index.",
                14, MUTED)

    round_rect(s, Inches(6.95), Inches(2.0), Inches(5.9), Inches(2.15), WHITE, LINE)
    rect(s, Inches(6.95), Inches(2.0), Inches(5.9), Inches(0.48), CYAN)
    add_textbox(s, Inches(7.15), Inches(2.05), Inches(5.5), Inches(0.4), "Semantic path  ·  MiniLM + Chroma", 16, WHITE, bold=True)
    add_textbox(s, Inches(7.15), Inches(2.6), Inches(5.5), Inches(1.4),
                "Catches paraphrases: “feeling hopeless” ≈ diagnostic symptoms.\nTop 50 nearest chunks in vector space.",
                14, MUTED)

    # fusion
    items = [
        ("RRF fusion", "Chunks found by BOTH paths rise. Reciprocal Rank Fusion, k = 60."),
        ("Cross-encoder rerank", "ms-marco MiniLM scores “does this passage actually answer this question?”"),
        ("Top evidence", "Return the best 5 chunks, with source + page, before any generation."),
    ]
    x = Inches(0.5)
    for title, body in items:
        round_rect(s, x, Inches(4.35), Inches(4.0), Inches(2.15), WHITE, LINE)
        add_textbox(s, x + Inches(0.2), Inches(4.5), Inches(3.6), Inches(0.55), title, 16, TEAL_DARK, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(5.1), Inches(3.6), Inches(1.2), body, 14, MUTED)
        x += Inches(4.15)

    footer(s, 4)
    notes(s, (
        "SPEAKER B — 45 seconds.\n\n"
        "When the user asks a question, we do not send it straight to the LLM.\n"
        "Arabic questions are translated to English for retrieval, then:\n"
        "1) BM25 keyword search — good for drug names and instrument names.\n"
        "2) Semantic search with all-MiniLM-L6-v2 in Chroma — good for paraphrases.\n"
        "3) Reciprocal Rank Fusion merges the two lists.\n"
        "4) A cross-encoder reranks the fused candidates.\n"
        "5) We keep the top 5 chunks as the evidence pack.\n\n"
        "Only then is generation allowed.\n\n"
        "بالعربي: السيستم بيدور في الدليل بطريقتين — كلمات مفتاحية ومعنى — وبعدين يدمجهم ويعمل rerank، وبعدين يدي الـ LLM الدليل بس."
    ))


def slide_hallucination(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.22), "04  ·  Hallucination prevention & citations   ·   MOST IMPORTANT")
    add_textbox(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.45),
                "The model is not allowed to invent. If it tries, we strip it.",
                24, NAVY, bold=True)

    guards = [
        (TEAL, "1  Grounded prompt", "Temperature 0. Every claim must use a provided [n]. No outside knowledge. If passages are weak, say so."),
        (CYAN, "2  Citation firewall", "We delete any [n] that was not in the retrieved set. Invented page numbers never reach the user."),
        (AMBER, "3  Relevance gate", "Cross-encoder score must clear 2.0. Below that: refuse. Weak but on-topic: “insufficient evidence”."),
        (RED, "4  Out-of-scope lock", "Off-topic questions (capital cities, recipes, code) return no answer. Confidence = 0."),
        (GREEN, "5  Objective confidence", "Score comes from retrieval + citation coverage — never from the LLM saying “I am sure”."),
        (TEAL_DARK, "6  Page-level proof", "UI shows PDF name, page, chunk, relevance, and a link to the official guideline."),
    ]
    positions = [
        (Inches(0.45), Inches(1.1)),
        (Inches(4.55), Inches(1.1)),
        (Inches(8.65), Inches(1.1)),
        (Inches(0.45), Inches(3.55)),
        (Inches(4.55), Inches(3.55)),
        (Inches(8.65), Inches(3.55)),
    ]
    for (fill, title, body), (x, y) in zip(guards, positions):
        round_rect(s, x, y, Inches(3.95), Inches(2.25), WHITE, LINE)
        rect(s, x, y, Inches(3.95), Inches(0.08), fill)
        add_textbox(s, x + Inches(0.18), y + Inches(0.22), Inches(3.6), Inches(0.55), title, 15, SLATE, bold=True)
        add_textbox(s, x + Inches(0.18), y + Inches(0.78), Inches(3.6), Inches(1.3), body, 13, MUTED)

    footer(s, 5)
    notes(s, (
        "SPEAKER B — 70–80 seconds. This is the slide the jury asked for. Slow down.\n\n"
        "Six locks, in order:\n"
        "1. Prompt: answer ONLY from numbered passages. Temperature 0.\n"
        "2. After generation we validate citations. Fake [7] when only [1]–[5] exist is deleted.\n"
        "3. Guardrail threshold: rerank score >= 2.0 or we do not generate.\n"
        "4. Off-topic → out_of_scope, empty sources, confidence 0.\n"
        "5. Mental-health question but weak chunks → insufficient evidence, not a guessed answer.\n"
        "6. Every surviving claim is shown with PDF, page, chunk, and official URL.\n\n"
        "If the LLM produces no valid citation, we fall back to an extractive snippet from the top chunk — still cited.\n\n"
        "النقطة الأهم: الـ LLM مش حر. لو evidenc مش كفاية أو السؤال برة النطاق، السيستم بيرفض. "
        "والإجابة المتدعمة بتتحط معاها رقم الصفحة والمصدر."
    ))


def slide_citations_example(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "04b  ·  What the jury will see on every answer")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12), Inches(0.45),
                "Evidence in. Citations out. Nothing in between is trusted.",
                24, NAVY, bold=True)

    # answer mock
    round_rect(s, Inches(0.5), Inches(1.2), Inches(7.7), Inches(5.35), WHITE, LINE)
    rect(s, Inches(0.5), Inches(1.2), Inches(7.7), Inches(0.5), TEAL)
    add_textbox(s, Inches(0.7), Inches(1.26), Inches(7.3), Inches(0.4),
                "Answer  ·  status: supported  ·  confidence 82%", 15, WHITE, bold=True)
    add_textbox(s, Inches(0.75), Inches(1.9), Inches(7.2), Inches(2.4),
                "NICE recommends a stepped-care approach for adults with depression, matching treatment intensity to severity. [1]\n\n"
                "For moderate depression this can include psychological therapy and, where indicated, antidepressant medication with monitoring. [1][2]\n\n"
                "Screening instruments such as PHQ-9 are used to assess severity in clinical settings. [3]",
                15, SLATE)

    add_textbox(s, Inches(0.75), Inches(4.45), Inches(7.2), Inches(0.35),
                "Citation firewall in action", 14, TEAL_DARK, bold=True)
    add_textbox(s, Inches(0.75), Inches(4.85), Inches(7.2), Inches(1.4),
                "LLM output:  “…as stated in [9] …”   →   [9] is not in the source index.\n"
                "We drop [9]. The sentence is either recast onto a valid ID, or the extractive fallback is used.\n"
                "The user never sees a fake guideline number.",
                14, MUTED)

    # source cards
    sources = [
        ("[1]", "Treatment Guidelines", "NG222  ·  page 48  ·  chunk 12", GREEN),
        ("[2]", "WHO Clinical Depression Guide", "WHO EMRO  ·  page 7  ·  chunk 3", TEAL),
        ("[3]", "Suicide Risk Assessment", "USPSTF adults  ·  page 4  ·  chunk 2", CYAN),
    ]
    y = 1.2
    for cid, name, meta, color in sources:
        round_rect(s, Inches(8.45), Inches(y), Inches(4.4), Inches(1.55), WHITE, LINE)
        oval(s, Inches(8.65), Inches(y + 0.22), Inches(0.5), Inches(0.5), color)
        add_textbox(s, Inches(8.65), Inches(y + 0.22), Inches(0.5), Inches(0.5),
                    cid, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, Inches(9.3), Inches(y + 0.22), Inches(3.3), Inches(0.4), name, 14, SLATE, bold=True)
        add_textbox(s, Inches(9.3), Inches(y + 0.62), Inches(3.3), Inches(0.7), meta + "\nOpen official guideline →", 13, MUTED)
        y += 1.7

    footer(s, 6)
    notes(s, (
        "SPEAKER B — 20 seconds, then pause so the jury can read.\n\n"
        "This is the product contract: every sentence has a [n], each [n] maps to a real PDF page, "
        "and the UI can open the official guideline.\n"
        "Invalid citations are removed in code, not by asking the model to behave."
    ))


def slide_evaluation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "05  ·  Evaluation  ·  then live demo")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12), Inches(0.45),
                "We measure retrieval, then we prove it on the Web UI.",
                24, NAVY, bold=True)

    metrics = [
        ("Recall@5", "Did the right guideline appear in the top 5 chunks?", "TEST_QUERIES\nsource-level recall"),
        ("Precision@k", "Of the top 1 / 3 / 5, how many are from the expected PDF?", "k = 1, 3, 5"),
        ("Guardrail", "Rerank threshold = 2.0\nOff-topic must return 0 hits", "out_of_scope = pass"),
        ("Confidence", "35% top score · 25% cited avg\n15% agreement · 15% grounding · 10% coverage", "never LLM self-score"),
    ]
    x = Inches(0.45)
    for title, body, foot in metrics:
        round_rect(s, x, Inches(1.2), Inches(3.05), Inches(3.15), WHITE, LINE)
        add_textbox(s, x + Inches(0.15), Inches(1.35), Inches(2.75), Inches(0.55), title, 18, TEAL_DARK, bold=True)
        add_textbox(s, x + Inches(0.15), Inches(1.95), Inches(2.75), Inches(1.35), body, 13, SLATE)
        add_textbox(s, x + Inches(0.15), Inches(3.4), Inches(2.75), Inches(0.7), foot, 12, MUTED, italic=True)
        x += Inches(3.2)

    # demo plan
    round_rect(s, Inches(0.45), Inches(4.55), Inches(12.45), Inches(2.05), NAVY)
    add_textbox(s, Inches(0.7), Inches(4.68), Inches(12), Inches(0.35),
                "Live demo — laptop operator, do not go back to slides", 16, LIGHT_TEAL, bold=True)
    add_textbox(s, Inches(0.7), Inches(5.1), Inches(12), Inches(1.3),
                "1) On-topic:  “What treatments does NICE recommend for moderate depression in adults?”   →  citations + pages\n"
                "2) Out of scope:  “What is the capital of France?”   →  refusal, confidence 0\n"
                "3) Arabic + voice:  “إيه أعراض الاكتئاب؟”   →  Egyptian Arabic answer with the same sources\n"
                "4) Screening:  PHQ-9 score  →  “Ask the guideline about this score”",
                14, WHITE)

    footer(s, 7)
    notes(s, (
        "SPEAKER B — 25 seconds on metrics, then SPEAKER A/operator opens the UI.\n\n"
        "We evaluate with Precision@k and Recall@k on held-out clinical queries, "
        "plus an out-of-scope check against the rerank guardrail (threshold 2.0).\n"
        "Confidence is a formula on retrieval signals, not a model opinion.\n\n"
        "Run `python main.py --eval` before the stage so you can quote the latest Precision@5 / Recall@5 if asked.\n\n"
        "Then: 'Let us show it live.' Stop talking. Click the browser.\n\n"
        "لو اللجنة سألت على الأرقام: احنا بنقيس Recall و Precision على مصادر الـ PDF المتوقعة، "
        "والـ guardrail بيتأكد إن السؤال اللي برة النطاق مايرجعش إجابة."
    ))


def slide_product(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, OFF_WHITE)
    kicker(s, Inches(0.5), Inches(0.28), "Bonus  ·  product, architecture, extras")
    add_textbox(s, Inches(0.5), Inches(0.58), Inches(12), Inches(0.45),
                "A real service — not a notebook demo.",
                24, NAVY, bold=True)

    extras = [
        ("Clean services", "ingest · retrieval · relevance · generation · citations · confidence · FastAPI  /api/chat"),
        ("Responsive UI", "React + Vite. Desktop sidebar, mobile nav, EN / المصري toggle, RTL."),
        ("Voice", "Speech-to-text questions and read-aloud answers (Arabic Egyptian + English)."),
        ("Screening trio", "PHQ-9, GAD-7, EPDS with crisis alert, then one-click guideline question."),
        ("Business case", "Clinics, university counseling, national helplines: faster, cited, bilingual support."),
        ("Career-ready", "Separation of concerns, cached indexes, evaluation harness, production-shaped API."),
    ]
    positions = [
        (Inches(0.45), Inches(1.2)),
        (Inches(4.55), Inches(1.2)),
        (Inches(8.65), Inches(1.2)),
        (Inches(0.45), Inches(3.55)),
        (Inches(4.55), Inches(3.55)),
        (Inches(8.65), Inches(3.55)),
    ]
    colors = [TEAL, CYAN, AMBER, GREEN, TEAL_DARK, SLATE]
    for (title, body), (x, y), c in zip(extras, positions, colors):
        round_rect(s, x, y, Inches(3.95), Inches(2.15), WHITE, LINE)
        rect(s, x, y, Inches(0.1), Inches(2.15), c)
        add_textbox(s, x + Inches(0.28), y + Inches(0.2), Inches(3.5), Inches(0.5), title, 16, SLATE, bold=True)
        add_textbox(s, x + Inches(0.28), y + Inches(0.75), Inches(3.5), Inches(1.2), body, 14, MUTED)

    footer(s, 8)
    notes(s, (
        "BACKUP SLIDE — skip if time is tight. Use in Q&A.\n\n"
        "Architecture is modular Python services + FastAPI + React. "
        "Bonus features the brief asked for: clean structure, responsive UI, voice, extra clinical screens.\n"
        "Business: this can sit in a clinic or helpline as a cited copilot — it does not replace a doctor."
    ))


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.18), H, TEAL)

    add_textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.9),
                "MindCare does not guess.", 36, WHITE, bold=True)
    add_textbox(s, Inches(0.7), Inches(2.45), Inches(12), Inches(1.1),
                "It retrieves official evidence, cites the page, and stays silent when the evidence is not there.",
                20, RGBColor(0xCB, 0xD5, 0xE1))
    add_textbox(s, Inches(0.7), Inches(3.6), Inches(12), Inches(0.5),
                "الدليل الأول. التخمين صفر.", 20, LIGHT_TEAL, font=FONT_AR)

    round_rect(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.35), TEAL_DEEP)
    add_textbox(s, Inches(0.95), Inches(4.65), Inches(11.4), Inches(1.05),
                "Thank you  ·  We are ready for questions\n"
                "Problem  →  Ingestion  →  Retrieval  →  Anti-hallucination  →  Evaluation & demo",
                16, WHITE, align=PP_ALIGN.CENTER)

    add_textbox(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.4),
                "ODC Hackathon  ·  Clinical RAG for depression  ·  Team MindCare",
                13, SOFT)

    notes(s, (
        "ALL — 10 seconds.\n"
        "Close on the contract: retrieve, cite, or refuse.\n"
        "Hand the room to the jury. Do not keep talking.\n"
        "Laptop stays on the chat screen for questions."
    ))


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_title(prs)
    slide_problem(prs)
    slide_ingestion(prs)
    slide_retrieval(prs)
    slide_hallucination(prs)
    slide_citations_example(prs)
    slide_evaluation(prs)
    slide_product(prs)
    slide_close(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
