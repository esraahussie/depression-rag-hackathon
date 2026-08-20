"""MindCare business-first pitch — convince on the idea, still cover the 5 jury steps."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
TEAL_DEEP = RGBColor(0x11, 0x4B, 0x4A)
NAVY = RGBColor(0x0B, 0x1F, 0x2A)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
SOFT = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF5, 0xF0)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
AMBER = RGBColor(0xC2, 0x7A, 0x1A)
RED = RGBColor(0xB9, 0x1C, 0x1C)
GREEN = RGBColor(0x04, 0x78, 0x57)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
LIGHT_TEAL = RGBColor(0xCC, 0xFB, 0xF1)
MIST = RGBColor(0xE6, 0xF4, 0xF1)

W = Inches(13.333)
H = Inches(7.5)
FONT = "Calibri"
FONT_AR = "Segoe UI"
OUT = Path(__file__).resolve().parent / "MindCare_Business_Pitch.pptx"


def set_run(run, text, size, color, bold=False, font=FONT, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rPr.append(ea)
    ea.set("typeface", FONT_AR)


def add_textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
                font=FONT, italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf._txBody.bodyPr.set(
            "anchor",
            {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor],
        )
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, size, color, bold, font, italic)
    return box


def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def round_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1.15)
    s.shadow.inherit = False
    return s


def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def chevron(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, n, total=8, dark=False):
    color = RGBColor(0x9C, 0xA3, 0xAF) if dark else SOFT
    add_textbox(
        slide, Inches(0.5), Inches(7.22), Inches(9.5), Inches(0.22),
        "MindCare  ·  A trusted clinical layer for depression care in Egypt",
        10, WHITE if dark else MUTED, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, Inches(11.2), Inches(7.22), Inches(1.7), Inches(0.22),
        f"{n}  /  {total}", 10, WHITE if dark else MUTED,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
    )


def kicker(slide, l, t, text, color=TEAL):
    oval(slide, l, t + Inches(0.07), Inches(0.12), Inches(0.12), color)
    add_textbox(
        slide, l + Inches(0.2), t, Inches(11), Inches(0.26),
        text.upper(), 11, color, bold=True,
    )


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)
    rect(s, 0, Inches(6.42), W, Inches(1.08), TEAL_DEEP)

    add_textbox(
        s, Inches(0.7), Inches(0.42), Inches(12), Inches(0.32),
        "THE IDEA  ·  NOT ANOTHER CHATBOT", 13, GOLD, bold=True,
    )
    add_textbox(
        s, Inches(0.7), Inches(1.15), Inches(12.2), Inches(1.7),
        "The world’s best depression\nguidelines already exist.",
        36, WHITE, bold=True,
    )
    add_textbox(
        s, Inches(0.7), Inches(3.05), Inches(11.8), Inches(1.05),
        "The people who need them — GPs, helpline staff, families — cannot use them.\nSo they ask ChatGPT. And ChatGPT guesses.",
        20, RGBColor(0xCB, 0xD5, 0xE1),
    )
    add_textbox(
        s, Inches(0.7), Inches(4.35), Inches(11.5), Inches(0.55),
        "MindCare puts official NICE, WHO and USPSTF evidence in their hands — in Egyptian Arabic — and stays silent when it is not sure.",
        16, LIGHT_TEAL,
    )
    add_textbox(
        s, Inches(0.7), Inches(5.1), Inches(11.5), Inches(0.4),
        "الفكرة: الدليل موجود. اللي محتاجه مش بيوصل له. احنا بنخلي الوصول آمن.",
        15, GOLD, font=FONT_AR,
    )
    add_textbox(
        s, Inches(0.7), Inches(6.62), Inches(9), Inches(0.65),
        "MindCare  ·  A deployable clinical copilot for depression  ·  ODC Hackathon",
        15, WHITE, bold=True,
    )
    add_textbox(
        s, Inches(10.3), Inches(6.62), Inches(2.6), Inches(0.65),
        "5 minutes  ·  then live", 14, LIGHT_TEAL, align=PP_ALIGN.RIGHT,
    )
    notes(
        s,
        "SPEAKER A — 15 seconds. Pause after the first sentence.\n\n"
        "The world’s best depression guidelines already exist. "
        "The people who need them cannot use them. So they ask ChatGPT — and ChatGPT guesses.\n\n"
        "MindCare is the trusted layer: official evidence, in Egyptian Arabic, "
        "and it stays silent when it is not sure.\n\n"
        "This is not a therapy bot. This is how Egypt can use AI in mental health without gambling on a hallucination.",
    )


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.26), "01  ·  The problem  ·  why this is a market, not a homework")
    add_textbox(
        s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
        "Egypt does not have a shortage of guidelines.\nIt has a shortage of people who can apply them.",
        24, NAVY, bold=True,
    )

    stats = [
        ("~1", "psychiatrist per\n100,000 people", "WHO / national workforce data"),
        ("~7%", "of Egyptian adults\nlive with depression", "Most common psychiatric diagnosis"),
        ("~4%", "of moderate–severe\ncases get treated in EMR", "WHO Eastern Mediterranean coverage"),
        ("minutes", "is all a GP has.\nNICE is hundreds of pages.", "The job to be done"),
    ]
    x = Inches(0.45)
    for num, label, foot in stats:
        round_rect(s, x, Inches(1.55), Inches(3.05), Inches(2.55), WHITE, LINE)
        add_textbox(s, x + Inches(0.12), Inches(1.7), Inches(2.8), Inches(0.7), num, 28, TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.12), Inches(2.4), Inches(2.8), Inches(0.9), label, 14, SLATE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.12), Inches(3.35), Inches(2.8), Inches(0.55), foot, 11, MUTED, italic=True, align=PP_ALIGN.CENTER)
        x += Inches(3.2)

    round_rect(s, Inches(0.45), Inches(4.3), Inches(6.15), Inches(2.5), WHITE, LINE)
    rect(s, Inches(0.45), Inches(4.3), Inches(0.1), Inches(2.5), RED)
    add_textbox(s, Inches(0.75), Inches(4.45), Inches(5.65), Inches(0.4), "What is already happening", 16, RED, bold=True)
    add_textbox(
        s, Inches(0.75), Inches(4.9), Inches(5.65), Inches(1.7),
        "Patients, students and even junior clinicians already ask general AI about symptoms, drugs and suicide risk.\n\n"
        "That model was not trained to refuse. It invents pages, mixes child and adult pathways, and sounds confident.",
        14, SLATE,
    )

    round_rect(s, Inches(6.8), Inches(4.3), Inches(6.05), Inches(2.5), WHITE, LINE)
    rect(s, Inches(6.8), Inches(4.3), Inches(0.1), Inches(2.5), TEAL)
    add_textbox(s, Inches(7.1), Inches(4.45), Inches(5.55), Inches(0.4), "The operational cost", 16, TEAL_DARK, bold=True)
    add_textbox(
        s, Inches(7.1), Inches(4.9), Inches(5.55), Inches(1.7),
        "Primary-care doctors, GSMHAT hotlines and university counselors carry the load — without a psychiatrist in the room.\n\n"
        "If they cannot open the guideline in seconds, they either delay care or guess.",
        14, SLATE,
    )
    footer(s, 2)
    notes(
        s,
        "SPEAKER A — 55 seconds.\n\n"
        "Open with the scene: a family physician has a few minutes. The NICE adult guideline is a book. "
        "Egypt has about one psychiatrist per 100,000 people. Depression is common. "
        "In the Eastern Mediterranean region, only a few percent of moderate-to-severe cases get treated.\n\n"
        "Medical value: earlier, safer decisions at the front line.\n"
        "Operational value: multiply scarce specialists instead of replacing them.\n\n"
        "The real competitor is already in people’s pockets: ChatGPT. "
        "That is the unsafe default we have to beat.\n\n"
        "بالعربي: المشكلة مش إن مفيش دليل. المشكلة إن الطبيب والمواطن مش بيوصلوا للدليل في الوقت الصح، "
        "فبيروحوا للذكاء الاصطناعي العادي — وهو بيخمّن.",
    )


def slide_who(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.26), "The business case  ·  who uses it, who buys it, why they care")
    add_textbox(
        s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.5),
        "We do not sell “AI”. We sell safer minutes at the front line.",
        24, NAVY, bold=True,
    )

    buyers = [
        (TEAL, "Primary care & clinics", "GP / family physician",
         "Ask the guideline during the visit. Screen with PHQ-9 / GAD-7 / EPDS. Refer with evidence, not memory."),
        (RGBColor(0x0E, 0x74, 0x90), "Public helplines", "GSMHAT · MoH · university lines",
         "Agents get a cited answer in Egyptian Arabic — or a clean refusal — instead of improvising under pressure."),
        (AMBER, "Telco & digital health", "Orange · platforms · insurers",
         "A wellness channel that will not invent medical advice. Brand-safe. Policy-safe. Ready to white-label."),
    ]
    x = Inches(0.45)
    for color, title, who, body in buyers:
        round_rect(s, x, Inches(1.25), Inches(4.05), Inches(3.35), WHITE, LINE)
        rect(s, x, Inches(1.25), Inches(4.05), Inches(0.1), color)
        add_textbox(s, x + Inches(0.22), Inches(1.5), Inches(3.6), Inches(0.45), title, 16, color, bold=True)
        add_textbox(s, x + Inches(0.22), Inches(1.95), Inches(3.6), Inches(0.35), who, 13, MUTED, italic=True)
        add_textbox(s, x + Inches(0.22), Inches(2.4), Inches(3.6), Inches(1.95), body, 14, SLATE)
        x += Inches(4.2)

    round_rect(s, Inches(0.45), Inches(4.8), Inches(12.45), Inches(2.0), NAVY)
    add_textbox(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.35), "Why this is investable now", 15, GOLD, bold=True)
    add_textbox(
        s, Inches(0.7), Inches(5.35), Inches(12), Inches(1.25),
        "Half of Egypt is under 25 and already lives on the phone. Digital mental-health demand is here. "
        "What the market is missing is not another chatbot — it is a product a ministry, a hospital legal team, "
        "or Orange can actually put their name on.\n"
        "MindCare’s contract: we cite the page, or we do not speak.",
        15, WHITE,
    )
    footer(s, 3)
    notes(
        s,
        "SPEAKER A — 40 seconds.\n\n"
        "Three buyers, one job: safer minutes at the front line.\n"
        "1) Clinics and GPs — decision support during the visit.\n"
        "2) Public helplines — GSMHAT-style agents who cannot wait for a specialist.\n"
        "3) Telcos and digital health — Orange is in this room. They need a wellness product that will not hallucinate.\n\n"
        "We are not competing with psychiatrists. We are multiplying them.\n"
        "The business is B2B / B2G license, per seat or per channel — not ads on sad users.",
    )


def slide_concept(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.26), "The concept  ·  one product, one promise")
    add_textbox(
        s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.55),
        "A clinical copilot that is allowed to know only what the guideline says.",
        22, NAVY, bold=True,
    )

    steps = [
        ("Ask", "In English or المصري.\nVoice or text."),
        ("Ground", "Search only inside\nofficial PDFs."),
        ("Prove", "Every sentence\ncarries a page."),
        ("Protect", "If evidence is weak\nor off-topic: silence."),
    ]
    x = Inches(0.45)
    for i, (title, body) in enumerate(steps):
        if i < 3:
            chevron(s, x, Inches(1.3), Inches(3.2), Inches(1.55), TEAL if i % 2 == 0 else TEAL_DARK)
            tw = Inches(2.7)
        else:
            round_rect(s, x, Inches(1.3), Inches(2.95), Inches(1.55), TEAL_DEEP)
            tw = Inches(2.7)
        add_textbox(s, x + Inches(0.18), Inches(1.4), tw, Inches(0.4), f"{i+1}. {title}", 18, WHITE, bold=True)
        add_textbox(s, x + Inches(0.18), Inches(1.85), tw, Inches(0.85), body, 14, LIGHT_TEAL)
        x += Inches(3.2)

    # before / after
    round_rect(s, Inches(0.45), Inches(3.15), Inches(6.1), Inches(3.55), WHITE, LINE)
    add_textbox(s, Inches(0.7), Inches(3.3), Inches(5.6), Inches(0.4), "Without MindCare", 16, RED, bold=True)
    items_l = [
        "Open a 200-page PDF and hope",
        "Or paste the question into ChatGPT",
        "No page. No audit trail. No Arabic pathway",
        "A wrong dose still sounds fluent",
        "Legal and clinical risk sits on the clinic",
    ]
    y = 3.8
    for item in items_l:
        oval(s, Inches(0.75), Inches(y + 0.08), Inches(0.14), Inches(0.14), RED)
        add_textbox(s, Inches(1.05), Inches(y), Inches(5.2), Inches(0.42), item, 14, SLATE)
        y += 0.5

    round_rect(s, Inches(6.8), Inches(3.15), Inches(6.05), Inches(3.55), WHITE, LINE)
    add_textbox(s, Inches(7.05), Inches(3.3), Inches(5.55), Inches(0.4), "With MindCare", 16, GREEN, bold=True)
    items_r = [
        "Answer in seconds, from NICE / WHO / USPSTF",
        "Citation: file, page, chunk, official link",
        "Egyptian Arabic + screening tools in the same app",
        "Off-topic and weak evidence are refused",
        "A log a hospital or ministry can defend",
    ]
    y = 3.8
    for item in items_r:
        oval(s, Inches(7.1), Inches(y + 0.08), Inches(0.14), Inches(0.14), GREEN)
        add_textbox(s, Inches(7.4), Inches(y), Inches(5.2), Inches(0.42), item, 14, SLATE)
        y += 0.5

    footer(s, 4)
    notes(
        s,
        "SPEAKER A — 40 seconds. This is the idea they must remember.\n\n"
        "Ask → Ground in official PDFs → Prove with a page → Protect by staying silent.\n\n"
        "Without us: ChatGPT or a closed PDF. With us: a cited, bilingual, refuse-when-unsure copilot "
        "plus PHQ-9 / GAD-7 / EPDS in the same product.\n\n"
        "If they remember one line: we cite the page, or we do not speak.",
    )


def slide_how(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.26), "02–03  ·  Data + retrieval  ·  told as a product, not a pipeline")
    add_textbox(
        s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.5),
        "We only ingest what a hospital would trust. Then we find the right page in seconds.",
        22, NAVY, bold=True,
    )

    round_rect(s, Inches(0.45), Inches(1.2), Inches(6.15), Inches(5.5), WHITE, LINE)
    rect(s, Inches(0.45), Inches(1.2), Inches(6.15), Inches(0.55), TEAL)
    add_textbox(s, Inches(0.7), Inches(1.28), Inches(5.7), Inches(0.4), "What goes in  ·  official evidence only", 16, WHITE, bold=True)
    add_textbox(
        s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(1.35),
        "12 guidelines from NICE, WHO and USPSTF — adults, children, chronic illness, suicide screening, perinatal depression. Not blogs. Not forums. Not the open web.",
        15, SLATE,
    )
    ingest = [
        ("Clean", "Remove headers, footers, noise — so the model never learns from a page number."),
        ("Chunk", "Bite-size passages (1000 characters, overlap) a clinician can actually cite."),
        ("Tag", "Every piece keeps its PDF name, page, and public URL — 848 evidence chunks."),
    ]
    y = 3.35
    for t, b in ingest:
        add_textbox(s, Inches(0.7), Inches(y), Inches(5.7), Inches(0.28), t, 14, TEAL_DARK, bold=True)
        add_textbox(s, Inches(0.7), Inches(y + 0.28), Inches(5.7), Inches(0.55), b, 13, MUTED)
        y += 0.95

    round_rect(s, Inches(6.8), Inches(1.2), Inches(6.05), Inches(5.5), WHITE, LINE)
    rect(s, Inches(6.8), Inches(1.2), Inches(6.05), Inches(0.55), RGBColor(0x0E, 0x74, 0x90))
    add_textbox(s, Inches(7.05), Inches(1.28), Inches(5.6), Inches(0.4), "What comes back  ·  the right page", 16, WHITE, bold=True)
    add_textbox(
        s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(1.2),
        "When someone asks in Arabic, we search in the language of the guideline — then answer in المصري.",
        15, SLATE,
    )
    retr = [
        ("Exact words", "Drug names, PHQ-9, NG222 — keyword search does not miss them."),
        ("Meaning", "“I feel empty” still finds diagnostic symptoms — semantic search."),
        ("Best evidence", "The two lists are fused and reranked. Only the top passages reach the answer."),
    ]
    y = 3.25
    for t, b in retr:
        add_textbox(s, Inches(7.05), Inches(y), Inches(5.6), Inches(0.28), t, 14, RGBColor(0x0E, 0x74, 0x90), bold=True)
        add_textbox(s, Inches(7.05), Inches(y + 0.28), Inches(5.6), Inches(0.55), b, 13, MUTED)
        y += 0.95

    footer(s, 5)
    notes(
        s,
        "SPEAKER B — 50 seconds. Cover jury points 2 and 3 without jargon.\n\n"
        "Data: 12 official PDFs only. Clean, chunk, tag with page + URL. 848 chunks.\n"
        "Retrieval: keyword + meaning, fused, reranked. Arabic question → English search → Arabic answer.\n"
        "Business line: a hospital legal team can name every document we are allowed to know.\n\n"
        "Do not say BM25 / RRF unless they ask. If they ask: hybrid BM25 + MiniLM, RRF, cross-encoder.",
    )


def slide_trust(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.22), "04  ·  Trust is the product  ·  hallucination prevention  ·  most important")
    add_textbox(
        s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.7),
        "A ministry will not buy fluency.\nThey will buy a system that is allowed to say “I don’t know.”",
        24, NAVY, bold=True,
    )

    locks = [
        (TEAL, "Grounded by design", "The model may only speak from the passages we retrieved. Temperature 0. No “general knowledge”."),
        (RGBColor(0x0E, 0x74, 0x90), "Citation firewall", "If it invents a [7] that was never retrieved, we delete it in code — before the user sees it."),
        (AMBER, "Evidence bar", "If the match is weak, we do not generate. Guardrail threshold on the reranker. No bluffing."),
        (RED, "Hard refuse", "Capital cities, recipes, homework: out of scope. Confidence = 0. Empty sources. Full stop."),
        (GREEN, "Honest uncertainty", "A real mental-health question with thin evidence → “not enough in the documents.” Still not a guess."),
        (TEAL_DEEP, "Audit trail", "PDF, page, chunk, official link, confidence from retrieval — never from the model saying it is sure."),
    ]
    positions = [
        (Inches(0.45), Inches(1.4)),
        (Inches(4.55), Inches(1.4)),
        (Inches(8.65), Inches(1.4)),
        (Inches(0.45), Inches(3.85)),
        (Inches(4.55), Inches(3.85)),
        (Inches(8.65), Inches(3.85)),
    ]
    for (color, title, body), (x, y) in zip(locks, positions):
        round_rect(s, x, y, Inches(3.95), Inches(2.2), WHITE, LINE)
        rect(s, x, y, Inches(3.95), Inches(0.08), color)
        add_textbox(s, x + Inches(0.2), y + Inches(0.22), Inches(3.55), Inches(0.5), title, 16, SLATE, bold=True)
        add_textbox(s, x + Inches(0.2), y + Inches(0.75), Inches(3.55), Inches(1.25), body, 13, MUTED)

    footer(s, 6)
    notes(
        s,
        "SPEAKER B — 70 seconds. Slow down. This is why they should believe the idea.\n\n"
        "Frame it as procurement, not ML: would Orange put their logo on ChatGPT medical answers? No.\n"
        "Would a hospital? No.\n"
        "Six locks: grounded prompt, citation firewall, score threshold 2.0, out-of-scope refuse, "
        "insufficient-evidence message, page-level audit trail.\n\n"
        "Closing line: fluency is cheap. Trust is the product.\n\n"
        "بالعربي: الوزارة والمستشفى مش هيشتروا كلام طليق. هيشتروا نظام مسموح له يقول مش عارف.",
    )


def slide_gtm_eval(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, CREAM)
    kicker(s, Inches(0.5), Inches(0.26), "05  ·  Proof, path to market, then the live product")
    add_textbox(
        s, Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.45),
        "We can measure it. We can ship it. We can show it in the next minute.",
        22, NAVY, bold=True,
    )

    # proof
    proofs = [
        ("Recall & Precision", "Did the right guideline land in the top 5? We score every test question."),
        ("Guardrail test", "Off-topic must return nothing. Threshold 2.0 is a product rule, not a hope."),
        ("Already a product", "React UI, FastAPI, EN + المصري, voice in/out, PHQ-9 · GAD-7 · EPDS."),
    ]
    x = Inches(0.45)
    for title, body in proofs:
        round_rect(s, x, Inches(1.15), Inches(4.05), Inches(1.85), WHITE, LINE)
        add_textbox(s, x + Inches(0.2), Inches(1.28), Inches(3.65), Inches(0.4), title, 15, TEAL_DARK, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(1.72), Inches(3.65), Inches(1.1), body, 13, SLATE)
        x += Inches(4.2)

    add_textbox(s, Inches(0.5), Inches(3.15), Inches(12), Inches(0.35), "Go-to-market — three doors, one engine", 16, NAVY, bold=True)

    doors = [
        ("Now", "University counseling & private clinics. Demo-ready. Bilingual. Screening in the same app."),
        ("Next", "Pilot with primary care / helplines. Arabic + citations = something a public body can defend."),
        ("Scale", "Telco wellness (Orange), then MENA: same engine, local official guidelines."),
    ]
    x = Inches(0.45)
    colors = [TEAL, RGBColor(0x0E, 0x74, 0x90), AMBER]
    for (title, body), color in zip(doors, colors):
        round_rect(s, x, Inches(3.55), Inches(4.05), Inches(1.55), WHITE, LINE)
        add_textbox(s, x + Inches(0.2), Inches(3.65), Inches(3.65), Inches(0.32), title, 14, color, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(4.0), Inches(3.65), Inches(0.95), body, 13, SLATE)
        x += Inches(4.2)

    round_rect(s, Inches(0.45), Inches(5.25), Inches(12.45), Inches(1.55), NAVY)
    add_textbox(s, Inches(0.7), Inches(5.35), Inches(12), Inches(0.32), "Live demo — this is the product, not a slide", 14, GOLD, bold=True)
    add_textbox(
        s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.95),
        "1)  “What does NICE recommend for moderate depression in adults?”  →  show the page.\n"
        "2)  “What is the capital of France?”  →  it must refuse. That refusal is the business.\n"
        "3)  “إيه أعراض الاكتئاب؟”  →  same evidence, Egyptian Arabic.",
        14, WHITE,
    )
    footer(s, 7)
    notes(
        s,
        "SPEAKER B — 35 seconds, then operator opens the UI.\n\n"
        "Proof: Precision/Recall@5, guardrail on out-of-scope, confidence from retrieval.\n"
        "GTM: clinics now → public helplines next → Orange / MENA scale.\n"
        "Revenue one-liner if asked: B2B/B2G license per seat or per channel. We do not monetize patients.\n\n"
        "Then stop. Demo.\n"
        "The refusal question is the business proof — do not skip it.",
    )


def slide_close(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, Inches(0.16), H, GOLD)

    add_textbox(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.35), "THE ASK", 13, GOLD, bold=True)
    add_textbox(
        s, Inches(0.7), Inches(1.25), Inches(12.1), Inches(1.4),
        "Do not fund another chatbot.\nFund the layer Egypt can actually deploy.",
        28, WHITE, bold=True,
    )
    add_textbox(
        s, Inches(0.7), Inches(2.85), Inches(12), Inches(0.7),
        "MindCare turns official depression evidence into a bilingual, cited, refuse-when-unsure copilot\nfor clinics, helplines, and digital health — starting today.",
        16, RGBColor(0xCB, 0xD5, 0xE1),
    )

    asks = [
        ("Ministry / Creativa", "A primary-care or helpline pilot. Policy-safe AI, not a demo on a laptop."),
        ("Orange / investors", "Distribution and scale. A wellness channel that will not hallucinate."),
        ("Clinics & universities", "Use it this semester. Screening + cited answers in one product."),
    ]
    x = Inches(0.7)
    for title, body in asks:
        round_rect(s, x, Inches(3.8), Inches(3.9), Inches(1.85), TEAL_DEEP)
        add_textbox(s, x + Inches(0.2), Inches(3.95), Inches(3.5), Inches(0.45), title, 14, GOLD, bold=True)
        add_textbox(s, x + Inches(0.2), Inches(4.45), Inches(3.5), Inches(1.0), body, 13, WHITE)
        x += Inches(4.05)

    add_textbox(
        s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.45),
        "We cite the page. Or we do not speak.     ·     الدليل… أو السكوت.",
        18, LIGHT_TEAL, bold=True,
    )
    add_textbox(
        s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.35),
        "Thank you — we will show it live, then we are ready for questions.",
        14, SOFT,
    )
    notes(
        s,
        "ALL — 10 seconds if you return after demo; otherwise skip and stay on the UI.\n\n"
        "Ask is specific: pilot with public care / helpline; distribution with Orange; "
        "adoption with clinics and universities.\n"
        "Last line, then silence: We cite the page, or we do not speak.",
    )


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_title(prs)
    slide_problem(prs)
    slide_who(prs)
    slide_concept(prs)
    slide_how(prs)
    slide_trust(prs)
    slide_gtm_eval(prs)
    slide_close(prs)
    prs.save(str(OUT))
    print(f"Saved {OUT}")


if __name__ == "__main__":
    main()
